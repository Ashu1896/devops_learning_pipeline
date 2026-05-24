import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from src.utils import logger

# Premium Design Theme Colors
DARK_BG = RGBColor(15, 23, 42)       # Slate 900
LIGHT_BG = RGBColor(248, 250, 252)   # Slate 50
TEXT_LIGHT = RGBColor(241, 245, 249) # Slate 100
TEXT_DARK = RGBColor(30, 41, 59)     # Slate 800
INDIGO = RGBColor(79, 70, 229)       # Indigo 600
TEAL = RGBColor(20, 184, 166)        # Teal 500
BORDER_COLOR = RGBColor(226, 232, 240) # Slate 200

def set_solid_background(slide, color):
    """
    Draws a full-slide rectangle to act as a custom color background.
    """
    left = top = Inches(0)
    width = Inches(13.333)
    height = Inches(7.5)
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    bg_shape.line.fill.background() # No border
    # Send to back is done by ordering (it must be drawn first!)
    return bg_shape

def add_header(slide, title_text, subtitle_text, dark_mode=False):
    """
    Creates a beautiful left-aligned header with Title and Subtitle.
    """
    header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(1.2))
    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
    
    # Title
    p_title = tf.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = "Arial"
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_LIGHT if dark_mode else TEXT_DARK
    
    # Subtitle
    p_sub = tf.add_paragraph()
    p_sub.text = subtitle_text
    p_sub.font.name = "Arial"
    p_sub.font.size = Pt(13)
    p_sub.font.color.rgb = TEAL if dark_mode else INDIGO
    p_sub.font.bold = True
    p_sub.space_before = Pt(4)

def format_bullet_points(tf, bullets, text_color, font_size=13):
    """
    Populates bullet points into a text frame with clean margins.
    """
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 or tf.paragraphs[0].text else tf.paragraphs[0]
        p.text = f"•   {bullet}"
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.space_after = Pt(8)
        p.line_spacing = 1.15

def generate_pptx(topic: str, slides_data: list, output_path: Path):
    """
    Generates a high-end programmatically designed PowerPoint presentation based on structured slides JSON.
    """
    logger.info(f"Generating PPTX for '{topic}' using python-pptx...")
    
    prs = Presentation()
    # Set to 16:9 Widescreen standard
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Use blank layouts (index 6 is usually completely blank)
    blank_layout = prs.slide_layouts[6]
    
    for slide_idx, slide_info in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        title = slide_info.get("title", f"DevOps - {topic}")
        subtitle = slide_info.get("subtitle", "")
        bullets = slide_info.get("bullets", [])
        takeaway = slide_info.get("takeaway", "")
        notes = slide_info.get("notes", "")
        
        # Add speaker notes
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
            
        # Determine layout type programmatically based on index
        if slide_idx == 0:
            # -------------------------------------------------------------
            # SLIDE 1: Premium Title Slide (Dark Theme)
            # -------------------------------------------------------------
            set_solid_background(slide, DARK_BG)
            
            # Draw beautiful side accent bar
            accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.1), Inches(3.2))
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = INDIGO
            accent_bar.line.fill.background()
            
            # Title & Subtitle text block
            text_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.5), Inches(3.2))
            tf = text_box.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
            
            p_top = tf.paragraphs[0]
            p_top.text = "6-MONTH DEVops MASTERCLASS"
            p_top.font.name = "Arial"
            p_top.font.size = Pt(13)
            p_top.font.bold = True
            p_top.font.color.rgb = TEAL
            p_top.space_after = Pt(10)
            
            p_title = tf.add_paragraph()
            p_title.text = title
            p_title.font.name = "Arial"
            p_title.font.size = Pt(44)
            p_title.font.bold = True
            p_title.font.color.rgb = TEXT_LIGHT
            
            p_sub = tf.add_paragraph()
            p_sub.text = subtitle if subtitle else "Enterprise Production Systems Engineering"
            p_sub.font.name = "Arial"
            p_sub.font.size = Pt(16)
            p_sub.font.color.rgb = RGBColor(148, 163, 184) # Slate 400
            p_sub.space_before = Pt(10)
            
            # Bottom metadata band
            meta_box = slide.shapes.add_textbox(Inches(1.2), Inches(6.2), Inches(10.5), Inches(0.6))
            p_meta = meta_box.text_frame.paragraphs[0]
            p_meta.text = "SRE & PLATFORM ENGINEERING GROUP  •  AUTOMATED PIPELINE GENERATION"
            p_meta.font.name = "Arial"
            p_meta.font.size = Pt(9.5)
            p_meta.font.bold = True
            p_meta.font.color.rgb = RGBColor(100, 116, 139) # Slate 500
            
        elif slide_idx == 4:
            # -------------------------------------------------------------
            # SLIDE 5: Strategic Roadmap & Timeline (Dark Theme)
            # -------------------------------------------------------------
            set_solid_background(slide, DARK_BG)
            add_header(slide, title, subtitle, dark_mode=True)
            
            # Draw a beautiful horizontal timeline line
            timeline_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.04))
            timeline_line.fill.solid()
            timeline_line.fill.fore_color.rgb = TEAL
            timeline_line.line.fill.background()
            
            # Add 3 horizontal milestone nodes
            node_width = Inches(3.5)
            positions = [Inches(0.8), Inches(4.9), Inches(9.0)]
            
            for idx, bullet in enumerate(bullets[:3]):
                left_pos = positions[idx]
                
                # Draw circular node point
                node_dot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_pos + Inches(1.5), Inches(3.68), Inches(0.28), Inches(0.28))
                node_dot.fill.solid()
                node_dot.fill.fore_color.rgb = INDIGO
                node_dot.line.color.rgb = TEAL
                node_dot.line.width = Pt(2)
                
                # Text box below the node
                m_box = slide.shapes.add_textbox(left_pos, Inches(4.2), node_width, Inches(2.2))
                mtf = m_box.text_frame
                mtf.word_wrap = True
                
                p_stage = mtf.paragraphs[0]
                p_stage.text = f"MILESTONE 0{idx+1}"
                p_stage.font.name = "Arial"
                p_stage.font.size = Pt(11)
                p_stage.font.bold = True
                p_stage.font.color.rgb = TEAL
                p_stage.space_after = Pt(6)
                p_stage.alignment = PP_ALIGN.CENTER
                
                p_text = mtf.add_paragraph()
                p_text.text = bullet
                p_text.font.name = "Arial"
                p_text.font.size = Pt(12)
                p_text.font.color.rgb = TEXT_LIGHT
                p_text.alignment = PP_ALIGN.CENTER
                
            # Key Takeaway Banner
            takeaway_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5))
            pt = takeaway_box.text_frame.paragraphs[0]
            pt.text = f"Goal: {takeaway}" if takeaway else "Strategic Goal: Expand architecture blueprints to scale globally."
            pt.font.name = "Arial"
            pt.font.size = Pt(11)
            pt.font.bold = True
            pt.font.color.rgb = TEAL
            pt.alignment = PP_ALIGN.CENTER
            
        else:
            # -------------------------------------------------------------
            # SLIDES 2, 3, 4: High-Fidelity 2-Column Split Dashboard (Light Theme)
            # -------------------------------------------------------------
            set_solid_background(slide, LIGHT_BG)
            add_header(slide, title, subtitle, dark_mode=False)
            
            # Left Column (Core Technical Points)
            left_col = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(5.6), Inches(4.2))
            ltf = left_col.text_frame
            ltf.word_wrap = True
            
            p_lh = ltf.paragraphs[0]
            p_lh.text = "ARCHITECTURAL BLUEPRINT & BLUEPRINTS"
            p_lh.font.name = "Arial"
            p_lh.font.size = Pt(11)
            p_lh.font.bold = True
            p_lh.font.color.rgb = INDIGO
            p_lh.space_after = Pt(10)
            
            left_bullets = bullets[:len(bullets)//2 + len(bullets)%2]
            format_bullet_points(ltf, left_bullets, TEXT_DARK)
            
            # Right Column (Beautiful dashboard Card Box with soft border & shadow)
            card_x = Inches(6.8)
            card_y = Inches(2.0)
            card_w = Inches(5.7)
            card_h = Inches(4.2)
            
            # Draw visual background card shape
            card_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, card_x, card_y, card_w, card_h)
            card_bg.fill.solid()
            card_bg.fill.fore_color.rgb = RGBColor(255, 255, 255) # Clean pure white card
            card_bg.line.color.rgb = BORDER_COLOR
            card_bg.line.width = Pt(1)
            
            # Accent bar on top of the card
            top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, card_x, card_y, card_w, Inches(0.12))
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = TEAL
            top_bar.line.fill.background()
            
            # Text Inside Right Column Card
            right_col = slide.shapes.add_textbox(card_x + Inches(0.3), card_y + Inches(0.3), card_w - Inches(0.6), card_h - Inches(0.6))
            rtf = right_col.text_frame
            rtf.word_wrap = True
            
            p_rh = rtf.paragraphs[0]
            p_rh.text = "ENTERPRISE PRO-TIPS & STRATEGIES"
            p_rh.font.name = "Arial"
            p_rh.font.size = Pt(11)
            p_rh.font.bold = True
            p_rh.font.color.rgb = TEAL
            p_rh.space_after = Pt(10)
            
            right_bullets = bullets[len(bullets)//2 + len(bullets)%2:]
            format_bullet_points(rtf, right_bullets, TEXT_DARK)
            
            # Key Takeaway footer inside the card
            if takeaway:
                p_take = rtf.add_paragraph()
                p_take.text = f"Key Takeaway: {takeaway}"
                p_take.font.name = "Arial"
                p_take.font.size = Pt(10.5)
                p_take.font.bold = True
                p_take.font.color.rgb = INDIGO
                p_take.space_before = Pt(14)
                
    # Ensure folder and save presentation
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    logger.info(f"Successfully generated beautiful presentation at: {output_path}")
